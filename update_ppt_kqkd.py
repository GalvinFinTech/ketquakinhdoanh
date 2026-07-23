import os
import json
import sys

try:
    from pptx import Presentation
    from pptx.util import Pt, Inches
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
except ImportError:
    print("❌ Thư viện 'python-pptx' chưa được cài đặt!")
    print("💡 Vui lòng cài đặt bằng lệnh: pip install python-pptx")
    sys.exit(1)

def fmt_vn(val, is_percent=False):
    """
    Định dạng số theo chuẩn Việt Nam:
    - Dấu chấm (.) phân cách hàng nghìn.
    - Dấu phẩy (,) phân cách phần thập phân (2 chữ số cố định).
    - Ví dụ: 1234.56 -> 1.234,56 | 0.9 -> 0,90%
    """
    if val is None or val == "nan":
        return "-"
    try:
        f_val = float(val)
    except (ValueError, TypeError):
        return "-"

    # Format floating number with standard US 2 decimal places first
    formatted = f"{f_val:,.2f}"
    
    # Swap commas and dots to enforce Vietnamese standard
    formatted = formatted.replace(",", "TEMP").replace(".", ",").replace("TEMP", ".")

    if is_percent:
        prefix = "+" if f_val > 0 else ""
        return f"{prefix}{formatted}%"
    
    return formatted

def update_powerpoint_kqkd():
    """
    Tự động hóa trình bày KQKD vào PowerPoint (Yuanta Việt Nam):
    1. Lấy dữ liệu Top Doanh nghiệp Vốn hóa Lớn & Vừa Tăng Trưởng Dương (>0%).
    2. Sử dụng 1 BẢNG ĐƠN DUY NHẤT gồm 11 CỘT CHỈ TIÊU TÀI CHÍNH CHI TIẾT.
    3. Định dạng số chuẩn Việt Nam (2 chữ số thập phân, phân cách phẩy thập phân).
    4. Cỡ chữ: Summary 11pt | Table 9pt.
    5. Đảm bảo an toàn cấu trúc XML, ép kiểu integer cho kích thước PPT.
    """
    script_dir = os.path.dirname(os.path.abspath(__file__))
    json_path = os.path.join(script_dir, "summary_report.json")
    
    # Đường dẫn file PowerPoint trên OneDrive
    ppt_path = r"C:\Users\vi.nguyen\OneDrive - CÔNG TY TRÁCH NHIỆM HỮU HẠN CHỨNG KHOÁN YUANTA VIỆT NAM\Retail Research - Tài liệu-MacBook Air của Nguyễn (85)-2\001 Morning Note\Morning Note_26.pptx"

    if not os.path.exists(json_path):
        print(f"❌ LỖI: Không tìm thấy file dữ liệu '{json_path}'.")
        print("💡 Vui lòng chạy file 'extract_kqkd.py' trước để sinh file summary_report.json.")
        return False

    if not os.path.exists(ppt_path):
        print(f"❌ LỖI: Không tìm thấy file PowerPoint tại:\n{ppt_path}")
        return False

    print("📊 Đang đọc dữ liệu từ summary_report.json...")
    try:
        with open(json_path, 'r', encoding='utf-8') as f:
            data = json.load(f)
    except Exception as e:
        print(f"❌ Lỗi khi đọc file JSON: {e}")
        return False

    exec_summary = data.get("executive_summary", {})
    top_data = data.get("top_30_large_mid_growth", data.get("top_30_bigcap_growth", []))

    pub_count = exec_summary.get("published_count", 0)
    tot_comp = exec_summary.get("total_companies", 0)
    mkt_cap_pct = exec_summary.get("published_market_cap_percentage", 0.0)
    prof_comp = exec_summary.get("profitable_companies", 0)
    loss_comp = exec_summary.get("loss_companies", 0)
    prof_ratio = exec_summary.get("profit_ratio_percentage", 0.0)

    print(f"📂 Đang mở file PowerPoint: Morning Note_26.pptx (Dữ liệu có {len(top_data)} DN)...")
    try:
        prs = Presentation(ppt_path)
    except Exception as e:
        print(f"❌ Không thể mở file PowerPoint. Hãy đảm bảo bạn đã ĐÓNG hoàn toàn PowerPoint!")
        print(f"Chi tiết lỗi: {e}")
        return False

    target_shape = None
    target_slide = None
    target_slide_idx = -1

    for slide_idx, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            if shape.name.strip().lower() == "kqkd":
                target_shape = shape
                target_slide = slide
                target_slide_idx = slide_idx + 1
                print(f"🎯 Đã tìm thấy SelectionPane 'kqkd' tại Slide số {target_slide_idx}!")
                break
        if target_shape:
            break

    if not target_shape:
        print("❌ LỖI: Không tìm thấy Shape/Frame nào có tên 'kqkd' trong Selection Pane.")
        return False

    pane_left = int(target_shape.left)
    pane_top = int(target_shape.top)
    pane_width = int(target_shape.width)
    pane_height = int(target_shape.height)

    # Yuanta Brand Color Palette
    COLOR_BLUE = RGBColor(0, 59, 112)       # Yuanta Primary Blue (#003B70)
    COLOR_DARK = RGBColor(30, 41, 59)       # Slate Dark Text (#1E293B)
    COLOR_BG_ALT = RGBColor(240, 244, 248)   # Light Zebra Row Background
    COLOR_WHITE = RGBColor(255, 255, 255)
    COLOR_GREEN = RGBColor(5, 150, 105)     # Positive Growth Green (#059669)
    COLOR_RED = RGBColor(220, 38, 38)       # Negative Growth Red (#DC2626)

    tf = target_shape.text_frame
    tf.word_wrap = True
    tf.margin_left = int(Inches(0.02))
    tf.margin_right = int(Inches(0.02))
    tf.margin_top = int(Inches(0.02))
    tf.margin_bottom = int(Inches(0.02))
    tf.clear()

    # Line 1: Header Title
    p0 = tf.paragraphs[0]
    p0.text = "TỔNG QUAN KẾT QUẢ KINH DOANH Q2"
    p0.font.bold = True
    p0.font.size = Pt(11)
    p0.font.name = "Segoe UI"
    p0.font.color.rgb = COLOR_BLUE
    p0.space_after = Pt(2)

    # Line 2: Published progress
    p1 = tf.add_paragraph()
    p1.font.size = Pt(11)
    p1.font.name = "Segoe UI"
    p1.space_after = Pt(1)
    
    run1_1 = p1.add_run()
    run1_1.text = "Tiến độ công bố: "
    run1_1.font.bold = True
    run1_1.font.color.rgb = COLOR_DARK
    
    run1_2 = p1.add_run()
    run1_2.text = f"{fmt_vn(pub_count).split(',')[0]}/{fmt_vn(tot_comp).split(',')[0]} doanh nghiệp "
    run1_2.font.color.rgb = COLOR_DARK

    run1_3 = p1.add_run()
    run1_3.text = f"(Chiếm {fmt_vn(mkt_cap_pct)}% tổng vốn hóa toàn thị trường)."
    run1_3.font.bold = True
    run1_3.font.color.rgb = COLOR_BLUE

    # Line 3: Profitability ratio
    p2 = tf.add_paragraph()
    p2.font.size = Pt(11)
    p2.font.name = "Segoe UI"
    p2.space_after = Pt(3)
    
    run2_1 = p2.add_run()
    run2_1.text = "Tỷ lệ có lãi: "
    run2_1.font.bold = True
    run2_1.font.color.rgb = COLOR_DARK

    run2_2 = p2.add_run()
    run2_2.text = f"{fmt_vn(prof_ratio)}% "
    run2_2.font.bold = True
    run2_2.font.color.rgb = COLOR_GREEN if prof_ratio >= 50 else COLOR_RED

    run2_3 = p2.add_run()
    run2_3.text = f"({fmt_vn(prof_comp).split(',')[0]} DN Lãi / {fmt_vn(loss_comp).split(',')[0]} DN Lỗ)."
    run2_3.font.color.rgb = COLOR_DARK

    # Line 4: Table Header Title
    actual_count = len(top_data)
    p3 = tf.add_paragraph()
    p3.text = f"TOP {actual_count} DOANH NGHIỆP VỐN HÓA LỚN VÀ VỪA TĂNG TRƯỜNG LNST Q2 MẠNH NHẤT:"
    p3.font.bold = True
    p3.font.size = Pt(11)
    p3.font.name = "Segoe UI"
    p3.font.color.rgb = COLOR_BLUE
    p3.space_after = Pt(2)

    summary_height = int(Inches(1.05))
    table_top = int(pane_top + summary_height)
    table_height = int(max(pane_height - summary_height, Inches(1.0)))
    table_width = int(pane_width)

    rows_needed = actual_count + 1
    cols_needed = 11

    headers = [
        "STT", "Mã CP", "Vốn hóa\n(tỷ)", 
        "DTT Q2\n(tỷ)", "TT DTT\nYoY", "LNST Q2\n(tỷ)", "TT LNST\nYoY",
        "DTT 6T\n(tỷ)", "TT DTT\n6M", "LNST 6T\n(tỷ)", "TT LNST\n6M"
    ]
    alignments = [
        PP_ALIGN.CENTER, PP_ALIGN.CENTER, PP_ALIGN.RIGHT,
        PP_ALIGN.RIGHT, PP_ALIGN.RIGHT, PP_ALIGN.RIGHT, PP_ALIGN.RIGHT,
        PP_ALIGN.RIGHT, PP_ALIGN.RIGHT, PP_ALIGN.RIGHT, PP_ALIGN.RIGHT
    ]

    col_ratios = [0.04, 0.08, 0.11, 0.09, 0.09, 0.09, 0.09, 0.09, 0.09, 0.11, 0.12]

    existing_table_shape = None
    for s in target_slide.shapes:
        if s.name == "kqkd_native_table_single" and s.has_table:
            existing_table_shape = s
            break

    if existing_table_shape and len(existing_table_shape.table.rows) != rows_needed:
        target_slide.shapes._spTree.remove(existing_table_shape._element)
        existing_table_shape = None

    if existing_table_shape:
        table_shape = existing_table_shape
        table = table_shape.table
        table_shape.left = int(pane_left)
        table_shape.top = int(table_top)
        table_shape.width = int(table_width)
        table_shape.height = int(table_height)
    else:
        table_shape = target_slide.shapes.add_table(rows_needed, cols_needed, int(pane_left), int(table_top), int(table_width), int(table_height))
        table_shape.name = "kqkd_native_table_single"
        table = table_shape.table

    for c_idx, ratio in enumerate(col_ratios):
        table.columns[c_idx].width = int(table_width * ratio)

    for col_idx, header_text in enumerate(headers):
        cell = table.cell(0, col_idx)
        cell.fill.solid()
        cell.fill.fore_color.rgb = COLOR_BLUE
        cell.margin_left = int(Inches(0.01))
        cell.margin_right = int(Inches(0.01))
        cell.margin_top = int(Inches(0.01))
        cell.margin_bottom = int(Inches(0.01))

        cell.text_frame.clear()
        p = cell.text_frame.paragraphs[0]
        p.text = header_text
        p.alignment = alignments[col_idx]
        p.font.bold = True
        p.font.size = Pt(9)
        p.font.name = "Segoe UI"
        p.font.color.rgb = COLOR_WHITE

    for row_idx, item in enumerate(top_data, 1):
        stt = row_idx
        macp = str(item.get("MaCP", "-")).strip()
        vh = item.get("VonHoa")
        dtt_q2 = item.get("DTT_Q2")
        tt_dtt_q_yoy = item.get("TT_DTT_Q_YoY")
        lnst_q2 = item.get("LNST_Q2")
        tt_lnst_q_yoy = item.get("TT_LNST_Q_YoY")
        dtt_6m = item.get("DTT_6M")
        tt_dtt_6m_yoy = item.get("TT_DTT_6M_YoY")
        lnst_6m = item.get("LNST_6M")
        tt_lnst_6m_yoy = item.get("TT_LNST_6M_YoY")

        row_data = [
            str(stt),
            macp,
            fmt_vn(vh),
            fmt_vn(dtt_q2),
            fmt_vn(tt_dtt_q_yoy, is_percent=True),
            fmt_vn(lnst_q2),
            fmt_vn(tt_lnst_q_yoy, is_percent=True),
            fmt_vn(dtt_6m),
            fmt_vn(tt_dtt_6m_yoy, is_percent=True),
            fmt_vn(lnst_6m),
            fmt_vn(tt_lnst_6m_yoy, is_percent=True)
        ]
        
        raw_growth_vals = [
            None, None, None,
            None, tt_dtt_q_yoy, None, tt_lnst_q_yoy,
            None, tt_dtt_6m_yoy, None, tt_lnst_6m_yoy
        ]

        bg_color = COLOR_BG_ALT if row_idx % 2 == 1 else COLOR_WHITE

        for col_idx, val in enumerate(row_data):
            cell = table.cell(row_idx, col_idx)
            cell.fill.solid()
            cell.fill.fore_color.rgb = bg_color
            cell.margin_left = int(Inches(0.01))
            cell.margin_right = int(Inches(0.01))
            cell.margin_top = int(Inches(0.005))
            cell.margin_bottom = int(Inches(0.005))

            cell.text_frame.clear()
            p = cell.text_frame.paragraphs[0]
            p.text = val
            p.alignment = alignments[col_idx]
            p.font.name = "Segoe UI"
            p.font.size = Pt(9)

            if col_idx == 1:
                p.font.bold = True
                p.font.color.rgb = COLOR_BLUE
            elif raw_growth_vals[col_idx] is not None:
                p.font.bold = True
                g_val = raw_growth_vals[col_idx]
                p.font.color.rgb = COLOR_GREEN if g_val > 0 else (COLOR_RED if g_val < 0 else COLOR_DARK)
            else:
                p.font.color.rgb = COLOR_DARK

    print("💾 Đang lưu lại thay đổi vào file PowerPoint...")
    try:
        prs.save(ppt_path)
        print(f"🎉 CẬP NHẬT THÀNH CÔNG! Bảng 11 Cột Chi Tiết ({actual_count} DN) đã trình bày hoàn hảo vào Slide số {target_slide_idx}.")
        return True
    except Exception as e:
        print(f"❌ KHÔNG THỂ LƯU FILE: Hãy đảm bảo bạn đã ĐÓNG file 'Morning Note_26.pptx' trước khi chạy script!")
        print(f"Chi tiết lỗi: {e}")
        return False

if __name__ == "__main__":
    print("=" * 70)
    print("🚀 CHƯƠNG TRÌNH TRÌNH BÀY BẢNG 11 CỘT KQKD VÀO POWERPOINT (YUANTA VIỆT NAM)")
    print("=" * 70)
    update_powerpoint_kqkd()